"""
PitchPal v2 - Pitch Deck Analyzer

Accepts PDF or PPTX uploads, extracts text for the ReAct agent,
and uses Gemini Vision to analyze design quality + narrative flow.

PDF  → PyMuPDF renders slides as PNG images → full visual analysis
PPTX → python-pptx extracts text          → content/narrative analysis only
"""

import io
import json
import logging
import re
from typing import Optional
import base64

logger = logging.getLogger(__name__)

# Max slides to send to Gemini Vision (to keep cost low)
MAX_VISION_SLIDES = 8
# Max file size: 20 MB
MAX_FILE_SIZE = 20 * 1024 * 1024


class DeckAnalysisError(Exception):
    pass


class DeckAnalyzer:
    """
    Analyzes pitch deck files (PDF/PPTX) using:
    1. PyMuPDF — PDF text extraction + slide rendering
    2. python-pptx — PPTX text extraction
    3. Gemini Vision — design quality + narrative analysis
    """

    def __init__(self, gemini_api_key: str, gemini_model: str = "gemini-2.5-flash"):
        self.api_key = gemini_api_key
        self.model_name = gemini_model

    # ─────────────────────────────────────────────────────────
    # Public API
    # ─────────────────────────────────────────────────────────

    async def analyze(self, file_bytes: bytes, filename: str) -> dict:
        """
        Main entry point.
        Returns a dict with keys:
          startup_name, extracted_text, slide_count, file_format, deck_quality
        """
        if len(file_bytes) > MAX_FILE_SIZE:
            raise DeckAnalysisError("File too large. Maximum size is 20 MB.")

        fname_lower = filename.lower()
        if fname_lower.endswith(".pdf"):
            return await self._analyze_pdf(file_bytes)
        elif fname_lower.endswith((".pptx", ".ppt")):
            return await self._analyze_pptx(file_bytes)
        else:
            raise DeckAnalysisError("Unsupported file format. Please upload a PDF or PPTX file.")

    # ─────────────────────────────────────────────────────────
    # PDF Analysis (text + images)
    # ─────────────────────────────────────────────────────────

    async def _analyze_pdf(self, file_bytes: bytes) -> dict:
        try:
            import fitz  # PyMuPDF
        except ImportError:
            raise DeckAnalysisError("PyMuPDF is not installed. Run: pip install PyMuPDF")

        doc = fitz.open(stream=file_bytes, filetype="pdf")
        try:
            slide_count = len(doc)
            extracted_text = self._extract_pdf_text(doc)
            slide_images = self._render_pdf_slides(doc, max_slides=MAX_VISION_SLIDES)
        finally:
            doc.close()

        # If PDF has no text layer (image-based deck), use Gemini Vision to OCR
        if extracted_text == "No text content found in PDF." and slide_images:
            logger.info("PDF has no text layer — using Gemini Vision to extract text from slide images")
            extracted_text = await self._ocr_slides_with_gemini(slide_images, slide_count)

        # Run Gemini Vision analysis
        deck_quality = await self._gemini_vision_analysis(
            slide_images=slide_images,
            extracted_text=extracted_text,
            has_images=True,
            slide_count=slide_count,
        )

        return {
            "startup_name": deck_quality.get("startup_name", "Unknown"),
            "extracted_text": extracted_text,
            "slide_count": slide_count,
            "file_format": "pdf",
            "deck_quality": deck_quality,
        }

    def _extract_pdf_text(self, doc) -> str:
        """Extract and clean text from all PDF pages."""
        pages = []
        for i, page in enumerate(doc, 1):
            text = page.get_text("text").strip()
            if text:
                pages.append(f"[Slide {i}]\n{text}")
        return "\n\n".join(pages) if pages else "No text content found in PDF."

    def _render_pdf_slides(self, doc, max_slides: int) -> list:
        """Render PDF pages as PNG bytes at 1.5x scale."""
        images = []
        for i, page in enumerate(doc):
            if i >= max_slides:
                break
            try:
                import fitz
                matrix = fitz.Matrix(1.5, 1.5)
                pix = page.get_pixmap(matrix=matrix, alpha=False)
                images.append(pix.tobytes("png"))
            except Exception as e:
                logger.warning(f"Could not render slide {i+1}: {e}")
        return images

    # ─────────────────────────────────────────────────────────
    # PPTX Analysis (text only)
    # ─────────────────────────────────────────────────────────

    async def _analyze_pptx(self, file_bytes: bytes) -> dict:
        try:
            from pptx import Presentation
        except ImportError:
            raise DeckAnalysisError("python-pptx is not installed. Run: pip install python-pptx")

        prs = Presentation(io.BytesIO(file_bytes))
        slide_count = len(prs.slides)

        # Extract text from all slides
        extracted_text = self._extract_pptx_text(prs)

        # Text-only Gemini analysis (no images for PPTX)
        deck_quality = await self._gemini_vision_analysis(
            slide_images=[],
            extracted_text=extracted_text,
            has_images=False,
            slide_count=slide_count,
        )

        return {
            "startup_name": deck_quality.get("startup_name", "Unknown"),
            "extracted_text": extracted_text,
            "slide_count": slide_count,
            "file_format": "pptx",
            "deck_quality": deck_quality,
        }

    def _extract_pptx_text(self, prs) -> str:
        """Extract text from all PPTX shapes."""
        slides_text = []
        for i, slide in enumerate(prs.slides, 1):
            slide_lines = []
            for shape in slide.shapes:
                if hasattr(shape, "text_frame"):
                    for para in shape.text_frame.paragraphs:
                        t = para.text.strip()
                        if t:
                            slide_lines.append(t)
            if slide_lines:
                slides_text.append(f"[Slide {i}]\n" + "\n".join(slide_lines))
        return "\n\n".join(slides_text) if slides_text else "No text content found in PPTX."

    # ─────────────────────────────────────────────────────────
    # Gemini Vision OCR (for image-based PDFs)
    # ─────────────────────────────────────────────────────────

    async def _ocr_slides_with_gemini(self, slide_images: list, slide_count: int) -> str:
        """Use Gemini Vision to extract text content from slide images when PDF has no text layer."""
        import google.generativeai as genai

        genai.configure(api_key=self.api_key)
        model = genai.GenerativeModel(
            model_name=self.model_name,
            generation_config={"temperature": 0.0, "max_output_tokens": 4096},
        )

        prompt = f"""You are reading a {slide_count}-slide pitch deck. Extract ALL visible text from each slide.

Format your output EXACTLY like this:
[Slide 1]
(all text from slide 1)

[Slide 2]
(all text from slide 2)

...and so on.

Rules:
- Extract every piece of visible text: titles, subtitles, bullet points, numbers, labels, captions.
- Preserve the reading order (top to bottom, left to right).
- Do NOT add commentary or interpretation — only extract what is written on the slides.
- If a slide has charts/graphs, extract any labels, axis titles, and data values visible."""

        parts: list = [prompt]
        for img_bytes in slide_images:
            parts.append({
                "mime_type": "image/png",
                "data": base64.b64encode(img_bytes).decode("utf-8"),
            })

        try:
            response = await model.generate_content_async(parts)
            extracted = response.text.strip()
            logger.info(f"Gemini OCR extracted {len(extracted)} chars from {len(slide_images)} slides")
            return extracted if extracted else "No text content found in PDF."
        except Exception as e:
            logger.error(f"Gemini OCR failed: {e}")
            return "No text content found in PDF."

    # ─────────────────────────────────────────────────────────
    # Gemini Vision Analysis
    # ─────────────────────────────────────────────────────────

    async def _gemini_vision_analysis(
        self,
        slide_images: list,
        extracted_text: str,
        has_images: bool,
        slide_count: int,
    ) -> dict:
        """Send slides + text to Gemini for structured quality analysis."""
        import google.generativeai as genai

        genai.configure(api_key=self.api_key)
        model = genai.GenerativeModel(
            model_name=self.model_name,
            generation_config={"temperature": 0.0, "max_output_tokens": 4096},
        )

        visual_note = (
            f"I am showing you {len(slide_images)} slide image(s) from a {slide_count}-slide pitch deck."
            if has_images
            else f"This is a {slide_count}-slide PowerPoint. I can only share the extracted text (no slide images available)."
        )

        if len(extracted_text) > 3000:
            logger.warning(f"Extracted text truncated from {len(extracted_text)} to 3000 chars for Gemini analysis")

        prompt = f"""You are an expert pitch deck designer and VC analyst with 15+ years of experience reviewing startup decks.

{visual_note}

Extracted slide text:
---
{extracted_text[:3000]}
---

Analyze this pitch deck and return ONLY a valid JSON object — no text before or after the JSON:
{{
  "startup_name": "Detected startup name from cover slide (or 'Unknown')",
  "design_score": 8.0,
  "narrative_score": 7.5,
  "data_viz_score": 7.0,
  "overall_deck_score": 7.5,
  "design_feedback": "One short sentence (max 100 chars) about visual design.",
  "narrative_feedback": "One short sentence (max 100 chars) about story flow.",
  "data_viz_feedback": "One short sentence (max 100 chars) about charts/data.",
  "strengths": ["Short strength 1", "Short strength 2"],
  "improvements": ["Short improvement 1", "Short improvement 2"]
}}

Scoring guide:
- design_score: Visual clarity, font consistency, color palette, whitespace use. Score 0 if no images available.
- narrative_score: Does the deck tell a compelling investor story? Is the flow logical?
- data_viz_score: Are charts clear and labeled? Are financial projections credible? Score 5 if no data visible.
- overall_deck_score: Weighted average (narrative 40%, design 35%, data_viz 25%)

CRITICAL RULES:
- Return ONLY the JSON object. No markdown code blocks, no explanation.
- Keep ALL feedback strings under 100 characters each.
- Maximum 2 items in strengths and improvements arrays.
- Keep it concise to avoid truncation."""

        # Build content parts: prompt text + optional slide images
        parts: list = [prompt]
        for img_bytes in slide_images:
            parts.append({
                "mime_type": "image/png",
                "data": base64.b64encode(img_bytes).decode("utf-8"),
            })

        try:
            response = await model.generate_content_async(parts)
            raw = response.text.strip()
            logger.info(f"Gemini deck analysis response (first 200): {raw[:200]}")
            return self._parse_deck_quality(raw, slide_count)
        except Exception as e:
            logger.error(f"Gemini Vision analysis failed: {e}")
            return self._fallback_deck_quality(slide_count)

    def _try_repair_json(self, raw: str) -> dict | None:
        """Attempt to parse and repair potentially truncated JSON."""
        # Strategy 1: direct parse
        try:
            return json.loads(raw)
        except json.JSONDecodeError:
            pass

        # Strategy 2: strip trailing commas
        try:
            fixed = re.sub(r",\s*([}\]])", r"\1", raw)
            return json.loads(fixed)
        except json.JSONDecodeError:
            pass

        # Strategy 3: extract largest {...} block
        try:
            start = raw.find("{")
            end = raw.rfind("}")
            if start != -1 and end > start:
                chunk = re.sub(r",\s*([}\]])", r"\1", raw[start:end + 1])
                return json.loads(chunk)
        except json.JSONDecodeError:
            pass

        # Strategy 4: close truncated JSON (missing brackets/braces)
        try:
            start = raw.find("{")
            if start != -1:
                partial = raw[start:]
                depth_brace = 0
                depth_bracket = 0
                in_str = False
                escape = False
                for ch in partial:
                    if escape:
                        escape = False
                        continue
                    if ch == '\\' and in_str:
                        escape = True
                        continue
                    if ch == '"' and not escape:
                        in_str = not in_str
                        continue
                    if in_str:
                        continue
                    if ch == '{':
                        depth_brace += 1
                    elif ch == '}':
                        depth_brace -= 1
                    elif ch == '[':
                        depth_bracket += 1
                    elif ch == ']':
                        depth_bracket -= 1

                if in_str:
                    partial += '"'
                partial += ']' * max(0, depth_bracket)
                partial += '}' * max(0, depth_brace)
                partial = re.sub(r",\s*([}\]])", r"\1", partial)
                return json.loads(partial)
        except (json.JSONDecodeError, Exception):
            pass

        return None

    def _parse_deck_quality(self, raw: str, slide_count: int) -> dict:
        """Parse Gemini's JSON response into a deck quality dict."""
        # Strip markdown code blocks if present
        raw = re.sub(r"^```(?:json)?\s*", "", raw, flags=re.MULTILINE)
        raw = re.sub(r"\s*```\s*$", "", raw, flags=re.MULTILINE)
        raw = raw.strip()

        data = self._try_repair_json(raw)
        if data is None:
            logger.warning("All JSON parse attempts failed for deck quality")
            return self._fallback_deck_quality(slide_count)

        return {
            "startup_name": str(data.get("startup_name", "Unknown")),
            "design_score": float(data.get("design_score", 6.0)),
            "narrative_score": float(data.get("narrative_score", 6.0)),
            "data_viz_score": float(data.get("data_viz_score", 5.0)),
            "overall_deck_score": float(data.get("overall_deck_score", 6.0)),
            "design_feedback": str(data.get("design_feedback", "No visual feedback available.")),
            "narrative_feedback": str(data.get("narrative_feedback", "Narrative analysis not available.")),
            "data_viz_feedback": str(data.get("data_viz_feedback", "No data visualization feedback.")),
            "strengths": list(data.get("strengths", [])),
            "improvements": list(data.get("improvements", [])),
            "analyzed_slides": slide_count,
        }

    def _fallback_deck_quality(self, slide_count: int) -> dict:
        """Return a default quality dict if Gemini analysis fails."""
        return {
            "startup_name": "Unknown",
            "design_score": 0.0,
            "narrative_score": 0.0,
            "data_viz_score": 0.0,
            "overall_deck_score": 0.0,
            "design_feedback": "Analysis unavailable — please try again.",
            "narrative_feedback": "Analysis unavailable — please try again.",
            "data_viz_feedback": "Analysis unavailable — please try again.",
            "strengths": [],
            "improvements": [],
            "analyzed_slides": slide_count,
        }
